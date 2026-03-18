"""Tests for file import handlers and FileHandlerRegistry."""

import io
import struct
import zlib

import pytest

from slack_agents import InputFile
from slack_agents.files import FileHandlerRegistry
from slack_agents.tools.base import FileImportToolException
from slack_agents.tools.file_importer import (
    Provider,
    _import_docx,
    _import_image,
    _import_pdf,
    _import_pptx,
    _import_text,
    _import_xlsx,
)


def _make_png():
    """Create a minimal valid PNG (1x1 red pixel)."""
    sig = b"\x89PNG\r\n\x1a\n"
    ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    ihdr_crc = zlib.crc32(b"IHDR" + ihdr_data) & 0xFFFFFFFF
    ihdr = struct.pack(">I", 13) + b"IHDR" + ihdr_data + struct.pack(">I", ihdr_crc)
    raw = zlib.compress(b"\x00\xff\x00\x00")
    idat_crc = zlib.crc32(b"IDAT" + raw) & 0xFFFFFFFF
    idat = struct.pack(">I", len(raw)) + b"IDAT" + raw + struct.pack(">I", idat_crc)
    iend_crc = zlib.crc32(b"IEND") & 0xFFFFFFFF
    iend = struct.pack(">I", 0) + b"IEND" + struct.pack(">I", iend_crc)
    return sig + ihdr + idat + iend


def _make_registry(allowed_functions=None):
    """Create a FileHandlerRegistry with the built-in Provider."""
    if allowed_functions is None:
        allowed_functions = [".*"]
    provider = Provider(allowed_functions=allowed_functions)
    return FileHandlerRegistry([provider])


# ---------------------------------------------------------------------------
# FileHandlerRegistry tests
# ---------------------------------------------------------------------------


def test_registry_routes_by_mime():
    registry = _make_registry()
    assert registry.can_handle("application/pdf")
    assert registry.can_handle("text/plain")
    assert registry.can_handle("image/png")
    assert not registry.can_handle("audio/mpeg")


def test_registry_supported_mimes():
    registry = _make_registry()
    mimes = registry.supported_mimes
    assert "application/pdf" in mimes
    assert "text/plain" in mimes
    assert "text/csv" in mimes
    assert "image/png" in mimes
    assert "image/jpeg" in mimes


async def test_registry_size_limit():
    registry = _make_registry()
    big_data = b"x" * 10_000_001
    result = await registry.process_file(big_data, "text/plain", "huge.txt", None, None)
    assert result is not None
    assert result["type"] == "text"
    assert "was not processed" in result["text"]
    assert "exceeds" in result["text"]
    assert "huge.txt" in result["text"]


async def test_registry_at_limit_is_ok():
    registry = _make_registry()
    data = b"A" * 10_000_000
    result = await registry.process_file(data, "text/plain", "exact.txt", None, None)
    assert result is not None
    assert "was not processed" not in result["text"]


async def test_registry_unhandled_mime():
    registry = _make_registry()
    result = await registry.process_file(b"data", "audio/mpeg", "song.mp3", None, None)
    assert result is None


def test_registry_empty_providers():
    registry = FileHandlerRegistry([])
    assert not registry.can_handle("text/plain")
    assert registry.supported_mimes == set()


def test_registry_allowed_functions_filter():
    """allowed_functions restricts which handlers are active."""
    registry = _make_registry(allowed_functions=["import_pdf", "import_text"])
    assert registry.can_handle("application/pdf")
    assert registry.can_handle("text/plain")
    assert not registry.can_handle("image/png")  # import_image not allowed


def test_registry_override_last_wins():
    """When multiple providers handle the same MIME, last one wins."""
    provider1 = Provider(allowed_functions=[".*"])
    provider2 = Provider(allowed_functions=[".*"])
    registry = FileHandlerRegistry([provider1, provider2])
    # Should not error — both handle text/plain, last wins
    assert registry.can_handle("text/plain")


# ---------------------------------------------------------------------------
# Text import handlers
# ---------------------------------------------------------------------------


def test_plain_text():
    f = InputFile(file_bytes=b"Hello, world!", mimetype="text/plain", filename="test.txt")
    result = _import_text(f)
    assert result is not None
    assert result["type"] == "text"
    assert "Hello, world!" in result["text"]


def test_csv():
    f = InputFile(
        file_bytes=b"name,age\nAlice,30\nBob,25", mimetype="text/csv", filename="data.csv"
    )
    result = _import_text(f)
    assert result is not None
    assert "Alice" in result["text"]
    assert "Bob" in result["text"]


def test_markdown():
    f = InputFile(
        file_bytes=b"# Title\n\nSome content.", mimetype="text/markdown", filename="readme.md"
    )
    result = _import_text(f)
    assert result is not None
    assert "# Title" in result["text"]


def test_utf8_with_errors():
    data = b"Hello \xff\xfe world"
    result = _import_text(InputFile(file_bytes=data, mimetype="text/plain", filename="bad.txt"))
    assert result is not None
    assert "Hello" in result["text"]
    assert "world" in result["text"]


# ---------------------------------------------------------------------------
# Image import handler
# ---------------------------------------------------------------------------


def test_image_import():
    png_bytes = _make_png()
    f = InputFile(file_bytes=png_bytes, mimetype="image/png", filename="test.png")
    result = _import_image(f)
    assert result is not None
    assert result["type"] == "image"
    assert result["source"]["type"] == "base64"
    assert result["source"]["media_type"] == "image/png"
    assert len(result["source"]["data"]) > 0


# ---------------------------------------------------------------------------
# PDF tests
# ---------------------------------------------------------------------------


def test_pdf_extraction_markdown():
    """Smoke test: pymupdf4llm produces markdown from a generated PDF."""
    from fpdf import FPDF

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=24)
    pdf.cell(text="Big Heading")
    pdf.ln()
    pdf.set_font("Helvetica", size=12)
    pdf.cell(text="Normal paragraph text here.")

    buf = pdf.output()
    f = InputFile(file_bytes=bytes(buf), mimetype="application/pdf", filename="test.pdf")
    result = _import_pdf(f)

    assert result is not None
    assert result["type"] == "text"
    assert "only text was extracted" in result["text"]
    assert "Big Heading" in result["text"]
    assert "Normal paragraph text here" in result["text"]


# ---------------------------------------------------------------------------
# XLSX tests
# ---------------------------------------------------------------------------


def test_xlsx_extraction_note():
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws["A1"] = "Revenue"
    ws["B1"] = 1000
    ws["A2"] = "Cost"
    ws["B2"] = 500

    buf = io.BytesIO()
    wb.save(buf)
    xlsx_bytes = buf.getvalue()

    mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    result = _import_xlsx(InputFile(file_bytes=xlsx_bytes, mimetype=mime, filename="test.xlsx"))

    assert result is not None
    assert "formulas appear as their" in result["text"]
    assert "pivot tables appear as their cached display values" in result["text"]
    assert "|" in result["text"]
    assert "---" in result["text"]
    assert "Revenue" in result["text"]
    assert "1000" in result["text"]


def test_xlsx_pipe_escaping():
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws["A1"] = "Header"
    ws["A2"] = "value|with|pipes"

    buf = io.BytesIO()
    wb.save(buf)
    xlsx_bytes = buf.getvalue()

    mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    result = _import_xlsx(InputFile(file_bytes=xlsx_bytes, mimetype=mime, filename="test.xlsx"))

    assert result is not None
    assert "value\\|with\\|pipes" in result["text"]


def test_xlsx_empty_sheet():
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Empty"

    buf = io.BytesIO()
    wb.save(buf)
    xlsx_bytes = buf.getvalue()

    mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    result = _import_xlsx(InputFile(file_bytes=xlsx_bytes, mimetype=mime, filename="test.xlsx"))

    assert result is not None
    assert "(empty)" in result["text"]


# ---------------------------------------------------------------------------
# DOCX tests
# ---------------------------------------------------------------------------


def test_docx_extraction_with_image():
    from docx import Document
    from docx.shared import Inches

    doc = Document()
    doc.add_paragraph("Hello from the document")

    png_bytes = _make_png()
    png_stream = io.BytesIO(png_bytes)
    doc.add_picture(png_stream, width=Inches(1))
    doc.add_paragraph("After the image")

    buf = io.BytesIO()
    doc.save(buf)
    docx_bytes = buf.getvalue()

    mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    result = _import_docx(InputFile(file_bytes=docx_bytes, mimetype=mime, filename="test.docx"))

    assert result is not None
    assert "Hello from the document" in result["text"]
    assert "After the image" in result["text"]
    assert "[IMAGE]" in result["text"]
    assert "Images found" in result["text"]


def test_docx_headings():
    from docx import Document

    doc = Document()
    doc.add_heading("Main Title", level=1)
    doc.add_heading("Sub Section", level=2)
    doc.add_paragraph("Body text")

    buf = io.BytesIO()
    doc.save(buf)
    docx_bytes = buf.getvalue()

    mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    result = _import_docx(InputFile(file_bytes=docx_bytes, mimetype=mime, filename="test.docx"))

    assert result is not None
    assert "# Main Title" in result["text"]
    assert "## Sub Section" in result["text"]
    assert "Body text" in result["text"]


def test_docx_bold_italic():
    from docx import Document

    doc = Document()
    p = doc.add_paragraph()
    p.add_run("normal ")
    bold_run = p.add_run("bold")
    bold_run.bold = True
    p.add_run(" ")
    italic_run = p.add_run("italic")
    italic_run.italic = True

    buf = io.BytesIO()
    doc.save(buf)
    docx_bytes = buf.getvalue()

    mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    result = _import_docx(InputFile(file_bytes=docx_bytes, mimetype=mime, filename="test.docx"))

    assert result is not None
    assert "**bold**" in result["text"]
    assert "*italic*" in result["text"]


def test_docx_table():
    from docx import Document

    doc = Document()
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Age"
    table.cell(1, 0).text = "Alice"
    table.cell(1, 1).text = "30"

    buf = io.BytesIO()
    doc.save(buf)
    docx_bytes = buf.getvalue()

    mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    result = _import_docx(InputFile(file_bytes=docx_bytes, mimetype=mime, filename="test.docx"))

    assert result is not None
    assert "|" in result["text"]
    assert "---" in result["text"]
    assert "Name" in result["text"]
    assert "Alice" in result["text"]


def test_docx_mixed_order():
    """Paragraph-table-paragraph ordering is preserved."""
    from docx import Document

    doc = Document()
    doc.add_paragraph("Before table")
    table = doc.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "Cell"
    doc.add_paragraph("After table")

    buf = io.BytesIO()
    doc.save(buf)
    docx_bytes = buf.getvalue()

    mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    result = _import_docx(InputFile(file_bytes=docx_bytes, mimetype=mime, filename="test.docx"))

    assert result is not None
    before_pos = result["text"].index("Before table")
    cell_pos = result["text"].index("Cell")
    after_pos = result["text"].index("After table")
    assert before_pos < cell_pos < after_pos


# ---------------------------------------------------------------------------
# PPTX tests
# ---------------------------------------------------------------------------


def test_pptx_extraction():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Introduction"
    slide.placeholders[1].text = "Welcome to the presentation"
    slide.notes_slide.notes_text_frame.text = "Speaker notes here"

    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    png_bytes = _make_png()
    png_stream = io.BytesIO(png_bytes)
    slide2.shapes.add_picture(png_stream, Inches(1), Inches(1), Inches(2), Inches(2))

    slide3 = prs.slides.add_slide(prs.slide_layouts[0])
    slide3.shapes.title.text = "Thank You"

    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    result = _import_pptx(InputFile(file_bytes=pptx_bytes, mimetype=mime, filename="test.pptx"))

    assert result is not None
    assert "[IMAGE]" in result["text"]
    assert "only text was extracted" in result["text"]
    assert "## Slide 1: Introduction" in result["text"]
    assert "Welcome to the presentation" in result["text"]
    assert "Notes: Speaker notes here" in result["text"]
    assert "## Slide 3: Thank You" in result["text"]


def test_pptx_table_extraction():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2)).table
    table.cell(0, 0).text = "Product"
    table.cell(0, 1).text = "Price"
    table.cell(1, 0).text = "Widget"
    table.cell(1, 1).text = "$10"

    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    result = _import_pptx(InputFile(file_bytes=pptx_bytes, mimetype=mime, filename="test.pptx"))

    assert result is not None
    assert "|" in result["text"]
    assert "---" in result["text"]
    assert "Product" in result["text"]
    assert "Widget" in result["text"]


def test_pptx_bold_text():
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Important"
    run.font.bold = True
    run.font.size = Pt(14)

    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    result = _import_pptx(InputFile(file_bytes=pptx_bytes, mimetype=mime, filename="test.pptx"))

    assert result is not None
    assert "**Important**" in result["text"]


# ---------------------------------------------------------------------------
# Provider tests
# ---------------------------------------------------------------------------


def test_provider_returns_all_tools():
    provider = Provider(allowed_functions=[".*"])
    tools = provider.tools
    names = {t["name"] for t in tools}
    assert "import_pdf" in names
    assert "import_docx" in names
    assert "import_xlsx" in names
    assert "import_pptx" in names
    assert "import_text" in names
    assert "import_image" in names


def test_provider_filters_by_allowed_functions():
    provider = Provider(allowed_functions=["import_pdf"])
    tools = provider.tools
    assert len(tools) == 1
    assert tools[0]["name"] == "import_pdf"


def test_input_file_provider_not_tool_provider():
    """BaseFileImporterProvider should not be a BaseToolProvider."""
    from slack_agents.tools.base import BaseToolProvider

    provider = Provider(allowed_functions=[".*"])
    assert not isinstance(provider, BaseToolProvider)


# ---------------------------------------------------------------------------
# Error-raise tests — importers raise FileImportToolException on bad data
# ---------------------------------------------------------------------------


def test_import_pdf_raises_on_bad_data():
    f = InputFile(file_bytes=b"not a pdf", mimetype="application/pdf", filename="bad.pdf")
    with pytest.raises(FileImportToolException, match="bad.pdf"):
        _import_pdf(f)


def test_import_docx_raises_on_bad_data():
    docx_mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    f = InputFile(file_bytes=b"not a docx", mimetype=docx_mime, filename="bad.docx")
    with pytest.raises(FileImportToolException, match="bad.docx"):
        _import_docx(f)


def test_import_xlsx_raises_on_bad_data():
    xlsx_mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    f = InputFile(file_bytes=b"not an xlsx", mimetype=xlsx_mime, filename="bad.xlsx")
    with pytest.raises(FileImportToolException, match="bad.xlsx"):
        _import_xlsx(f)


def test_import_pptx_raises_on_bad_data():
    pptx_mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    f = InputFile(file_bytes=b"not a pptx", mimetype=pptx_mime, filename="bad.pptx")
    with pytest.raises(FileImportToolException, match="bad.pptx"):
        _import_pptx(f)


async def test_provider_raises_on_unknown_handler():
    provider = Provider(allowed_functions=[".*"])
    f = InputFile(file_bytes=b"data", mimetype="text/plain", filename="test.txt")
    with pytest.raises(FileImportToolException, match="Unknown import handler"):
        await provider.call_tool("nonexistent_handler", f, None, None)


async def test_registry_returns_none_on_import_error():
    """FileHandlerRegistry.process_file returns None when import raises."""
    registry = _make_registry()
    result = await registry.process_file(b"not a pdf", "application/pdf", "bad.pdf", None, None)
    assert result is None
