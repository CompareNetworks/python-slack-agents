"""Tests for document generation tool provider."""

import io
import zipfile
from unittest.mock import MagicMock

import openpyxl
import pytest
from docx import Document

from slack_agents import UserConversationContext
from slack_agents.config import load_plugin
from slack_agents.tools.file_exporter import (
    _FONT_REGULAR,
    _TOOL_MANIFEST,
    _export_csv,
    _export_docx,
    _export_pdf,
    _export_pptx,
    _export_xlsx,
    _parse_body_lines,
)


class TestManifest:
    def test_all_entries_have_required_fields(self):
        for entry in _TOOL_MANIFEST:
            assert "name" in entry
            assert "description" in entry
            assert "input_schema" in entry
            assert "handler" in entry

    def test_handler_references_are_callable(self):
        for entry in _TOOL_MANIFEST:
            assert callable(entry["handler"])


class TestToolProvider:
    def test_load_provider(self):
        provider = load_plugin("slack_agents.tools.file_exporter", allowed_functions=[".*"])
        assert len(provider.tools) == 5
        names = {t["name"] for t in provider.tools}
        assert names == {"export_pdf", "export_docx", "export_xlsx", "export_csv", "export_pptx"}

    def test_allowed_functions_filter(self):
        provider = load_plugin(
            "slack_agents.tools.file_exporter",
            allowed_functions=["export_pdf", "export_docx", "export_csv"],
        )
        names = {t["name"] for t in provider.tools}
        assert "export_xlsx" not in names
        assert "export_pdf" in names
        assert "export_docx" in names

    def test_allowed_functions_regex(self):
        provider = load_plugin(
            "slack_agents.tools.file_exporter",
            allowed_functions=["export_pdf"],
        )
        assert len(provider.tools) == 1
        assert provider.tools[0]["name"] == "export_pdf"

    def test_allowed_functions_empty_allows_nothing(self):
        provider = load_plugin("slack_agents.tools.file_exporter", allowed_functions=[])
        assert provider.tools == []

    def test_missing_module_raises(self):
        with pytest.raises(ModuleNotFoundError):
            load_plugin("nonexistent.tools.fake", allowed_functions=[".*"])

    async def test_call_tool(self):
        provider = load_plugin("slack_agents.tools.file_exporter", allowed_functions=[".*"])
        ctx = UserConversationContext(
            user_id="U123",
            user_name="test",
            user_handle="test",
            channel_id="C001",
            channel_name="general",
            thread_id="1234.5678",
        )
        storage = MagicMock()
        result = await provider.call_tool(
            "export_pdf", {"title": "Test", "body": "Hello"}, ctx, storage
        )
        assert not result["is_error"]
        assert result["files"][0]["data"][:5] == b"%PDF-"

    async def test_call_unknown_tool(self):
        provider = load_plugin("slack_agents.tools.file_exporter", allowed_functions=[".*"])
        ctx = UserConversationContext(
            user_id="U123",
            user_name="test",
            user_handle="test",
            channel_id="C001",
            channel_name="general",
            thread_id="1234.5678",
        )
        storage = MagicMock()
        result = await provider.call_tool("nonexistent", {}, ctx, storage)
        assert result["is_error"]


class TestExportPdf:
    async def test_produces_valid_pdf(self):
        result = await _export_pdf({"title": "Test Report", "body": "Hello world"})
        assert not result["is_error"]
        assert len(result["files"]) == 1
        file_entry = result["files"][0]
        assert file_entry["filename"].endswith(".pdf")
        assert file_entry["mimeType"] == "application/pdf"
        assert file_entry["data"][:5] == b"%PDF-"

    async def test_handles_formatting(self):
        body = "# Heading 1\n## Heading 2\n- Bullet item\n**Bold text** here\nPlain text"
        result = await _export_pdf({"title": "Formatted", "body": body})
        assert not result["is_error"]
        assert result["files"][0]["data"][:5] == b"%PDF-"

    @pytest.mark.skipif(
        not _FONT_REGULAR.exists(),
        reason="DejaVu fonts not installed — Unicode chars require TTF fonts",
    )
    async def test_handles_unicode(self):
        body = (
            "Key points \u2013 overview\n"
            "- \u201csmart quotes\u201d and \u2018single\u2019\n"
            "- ellipsis\u2026 and bullets \u2022"
        )
        result = await _export_pdf({"title": "Unicode \u2014 Test", "body": body})
        assert not result["is_error"]
        assert result["files"][0]["data"][:5] == b"%PDF-"

    async def test_summary_message(self):
        result = await _export_pdf({"title": "My Doc", "body": "Content"})
        assert "My_Doc.pdf" in result["content"]
        assert "page" in result["content"]


class TestExportDocx:
    async def test_produces_valid_docx(self):
        result = await _export_docx({"title": "Test Doc", "body": "Hello world"})
        assert not result["is_error"]
        assert len(result["files"]) == 1
        file_entry = result["files"][0]
        assert file_entry["filename"].endswith(".docx")
        # DOCX is a ZIP file
        assert zipfile.is_zipfile(io.BytesIO(file_entry["data"]))

    async def test_handles_formatting(self):
        body = "# Heading 1\n## Heading 2\n- Bullet item\n**Bold text** here"
        result = await _export_docx({"title": "Formatted", "body": body})
        assert not result["is_error"]
        assert zipfile.is_zipfile(io.BytesIO(result["files"][0]["data"]))

    async def test_summary_message(self):
        result = await _export_docx({"title": "My Doc", "body": "Content"})
        assert "My_Doc.docx" in result["content"]


class TestExportXlsx:
    async def test_produces_valid_xlsx(self):
        result = await _export_xlsx(
            {
                "filename": "data",
                "sheets": [{"name": "Sheet1", "rows": [["Name", "Email"], ["Alice", "a@b.com"]]}],
            }
        )
        assert not result["is_error"]
        assert len(result["files"]) == 1
        file_entry = result["files"][0]
        assert file_entry["filename"].endswith(".xlsx")
        # Verify it opens with openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(file_entry["data"]))
        ws = wb["Sheet1"]
        assert ws.cell(1, 1).value == "Name"
        assert ws.cell(2, 1).value == "Alice"

    async def test_multiple_sheets(self):
        result = await _export_xlsx(
            {
                "filename": "multi",
                "sheets": [
                    {"name": "People", "rows": [["Name"], ["Bob"]]},
                    {"name": "Items", "rows": [["Item"], ["Widget"]]},
                ],
            }
        )
        assert not result["is_error"]
        wb = openpyxl.load_workbook(io.BytesIO(result["files"][0]["data"]))
        assert "People" in wb.sheetnames
        assert "Items" in wb.sheetnames

    async def test_summary_message(self):
        result = await _export_xlsx(
            {
                "filename": "report",
                "sheets": [{"name": "Data", "rows": [["A"], ["B"], ["C"]]}],
            }
        )
        assert "1 sheet" in result["content"]
        assert "3 rows" in result["content"]


class TestExportCsv:
    async def test_produces_valid_csv(self):
        result = await _export_csv(
            {
                "filename": "data",
                "rows": [["Name", "Email"], ["Alice", "a@b.com"]],
            }
        )
        assert not result["is_error"]
        assert len(result["files"]) == 1
        file_entry = result["files"][0]
        assert file_entry["filename"].endswith(".csv")
        assert file_entry["mimeType"] == "text/csv"
        content = file_entry["data"].decode("utf-8")
        assert "Name,Email" in content
        assert "Alice,a@b.com" in content

    async def test_handles_commas_and_quotes(self):
        result = await _export_csv(
            {
                "filename": "tricky",
                "rows": [["Name", "Note"], ["Alice", 'Said "hello", then left']],
            }
        )
        assert not result["is_error"]
        content = result["files"][0]["data"].decode("utf-8")
        # csv.writer should quote fields containing commas/quotes
        assert "Alice" in content
        assert "hello" in content

    async def test_summary_message(self):
        result = await _export_csv(
            {
                "filename": "report",
                "rows": [["A"], ["B"], ["C"]],
            }
        )
        assert "report.csv" in result["content"]
        assert "3 rows" in result["content"]


class TestExportPptx:
    async def test_produces_valid_pptx(self):
        result = await _export_pptx(
            {
                "title": "Test Presentation",
                "slides": [
                    {"title": "Slide One", "body": "Hello world"},
                    {"title": "Slide Two", "body": "- Bullet A\n- Bullet B"},
                ],
            }
        )
        assert not result["is_error"]
        assert len(result["files"]) == 1
        file_entry = result["files"][0]
        assert file_entry["filename"].endswith(".pptx")
        assert zipfile.is_zipfile(io.BytesIO(file_entry["data"]))

    async def test_handles_bullet_formatting(self):
        result = await _export_pptx(
            {
                "title": "Bullets",
                "slides": [{"title": "Key Points", "body": "- First\n- Second\n- Third"}],
            }
        )
        assert not result["is_error"]
        # Verify the PPTX contains the expected text
        from pptx import Presentation

        prs = Presentation(io.BytesIO(result["files"][0]["data"]))
        slide = prs.slides[0]
        body_text = slide.placeholders[1].text_frame.text
        assert "First" in body_text
        assert "Second" in body_text
        assert "Third" in body_text

    async def test_summary_message(self):
        result = await _export_pptx(
            {
                "title": "My Deck",
                "slides": [
                    {"title": "One"},
                    {"title": "Two"},
                    {"title": "Three"},
                ],
            }
        )
        assert "My_Deck.pptx" in result["content"]
        assert "3 slides" in result["content"]

    async def test_table_rendering(self):
        body = "| Name | Age |\n|---|---|\n| Alice | 30 |\n| Bob | 25 |"
        result = await _export_pptx(
            {"title": "Table Deck", "slides": [{"title": "Data", "body": body}]}
        )
        assert not result["is_error"]
        from pptx import Presentation

        prs = Presentation(io.BytesIO(result["files"][0]["data"]))
        slide = prs.slides[0]
        # Find the table shape
        table_shapes = [s for s in slide.shapes if s.has_table]
        assert len(table_shapes) == 1
        tbl = table_shapes[0].table
        assert len(tbl.rows) == 3  # header + 2 data rows
        assert tbl.cell(0, 0).text == "Name"
        assert tbl.cell(1, 0).text == "Alice"
        assert tbl.cell(2, 1).text == "25"

    async def test_bold_italic_text(self):
        body = "**bold text** and *italic text* and plain"
        result = await _export_pptx(
            {"title": "Rich Text", "slides": [{"title": "Styles", "body": body}]}
        )
        assert not result["is_error"]
        from pptx import Presentation

        prs = Presentation(io.BytesIO(result["files"][0]["data"]))
        slide = prs.slides[0]
        tf = slide.placeholders[1].text_frame
        runs = tf.paragraphs[0].runs
        # Find bold and italic runs
        bold_runs = [r for r in runs if r.font.bold]
        italic_runs = [r for r in runs if r.font.italic]
        assert any("bold text" in r.text for r in bold_runs)
        assert any("italic text" in r.text for r in italic_runs)


class TestParseBodyLines:
    def test_table_detection(self):
        body = "| Name | Age |\n|---|---|\n| Alice | 30 |"
        elements = list(_parse_body_lines(body))
        assert len(elements) == 1
        assert elements[0][0] == "table"
        rows = elements[0][1]
        assert len(rows) == 2  # header + 1 data row (separator skipped)
        assert rows[0] == ["Name", "Age"]
        assert rows[1] == ["Alice", "30"]

    def test_table_with_surrounding_text(self):
        body = "Intro text\n| A | B |\n|---|---|\n| 1 | 2 |\nAfter text"
        elements = list(_parse_body_lines(body))
        types = [e[0] for e in elements]
        assert types == ["paragraph", "table", "paragraph"]
        assert elements[0][1] == "Intro text"
        assert elements[1][1] == [["A", "B"], ["1", "2"]]
        assert elements[2][1] == "After text"

    def test_h3_heading(self):
        body = "### Sub-sub-heading"
        elements = list(_parse_body_lines(body))
        assert elements == [("h3", "Sub-sub-heading")]

    def test_numbered_list(self):
        body = "1. First item\n2. Second item\n3. Third item"
        elements = list(_parse_body_lines(body))
        assert all(e[0] == "numbered" for e in elements)
        assert elements[0][1] == "First item"
        assert elements[1][1] == "Second item"
        assert elements[2][1] == "Third item"

    def test_mixed_elements(self):
        body = "# Title\n## Section\n### Subsection\n- bullet\n1. numbered\nplain"
        elements = list(_parse_body_lines(body))
        types = [e[0] for e in elements]
        assert types == ["h1", "h2", "h3", "bullet", "numbered", "paragraph"]


class TestExportPdfTable:
    async def test_table_rendering(self):
        body = "Header text\n| Col A | Col B |\n|---|---|\n| X | Y |"
        result = await _export_pdf({"title": "Table PDF", "body": body})
        assert not result["is_error"]
        assert result["files"][0]["data"][:5] == b"%PDF-"


class TestExportDocxExtended:
    async def test_table_rendering(self):
        body = "| Name | Role |\n|---|---|\n| Alice | Dev |\n| Bob | PM |"
        result = await _export_docx({"title": "Table Doc", "body": body})
        assert not result["is_error"]
        doc = Document(io.BytesIO(result["files"][0]["data"]))
        assert len(doc.tables) == 1
        tbl = doc.tables[0]
        assert tbl.rows[0].cells[0].text == "Name"
        assert tbl.rows[0].cells[1].text == "Role"
        assert tbl.rows[1].cells[0].text == "Alice"
        assert tbl.rows[2].cells[1].text == "PM"

    async def test_italic_text(self):
        body = "This has *italic words* in it"
        result = await _export_docx({"title": "Italic Doc", "body": body})
        assert not result["is_error"]
        doc = Document(io.BytesIO(result["files"][0]["data"]))
        # Find paragraph with italic run (skip heading)
        for para in doc.paragraphs:
            for run in para.runs:
                if run.italic and "italic words" in run.text:
                    return
        pytest.fail("No italic run found with expected text")

    async def test_numbered_list(self):
        body = "1. First\n2. Second"
        result = await _export_docx({"title": "Numbered", "body": body})
        assert not result["is_error"]
        doc = Document(io.BytesIO(result["files"][0]["data"]))
        numbered_paras = [p for p in doc.paragraphs if p.style.name == "List Number"]
        assert len(numbered_paras) == 2

    async def test_h3_heading(self):
        body = "### My Subheading"
        result = await _export_docx({"title": "H3 Doc", "body": body})
        assert not result["is_error"]
        doc = Document(io.BytesIO(result["files"][0]["data"]))
        h3_paras = [p for p in doc.paragraphs if p.style.name == "Heading 3"]
        assert len(h3_paras) == 1
        assert h3_paras[0].text == "My Subheading"
