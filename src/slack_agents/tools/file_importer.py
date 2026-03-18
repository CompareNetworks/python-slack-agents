"""Built-in file input provider: document import (PDF, DOCX, XLSX, PPTX, text, images).

Exports a Provider class that subclasses BaseFileImporterProvider.
"""

import base64
import io
import logging

from slack_agents import InputFile, UserConversationContext
from slack_agents.storage.base import BaseStorageProvider
from slack_agents.tools.base import BaseFileImporterProvider, ContentBlock, FileImportToolException

logger = logging.getLogger(__name__)

MIME_DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
MIME_XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
MIME_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

_TEXT_ONLY_NOTE = (
    "(Note: only text was extracted; images are shown as [IMAGE] placeholders"
    " with alt text when available. Charts and visual formatting are not included.)"
)
_CELLS_ONLY_NOTE = (
    "(Note: only cell values were extracted — formulas appear as their"
    " last-calculated values (may be blank if never opened in Excel),"
    " and pivot tables appear as their cached display values."
    " Images, charts, and formatting are not included.)"
)


def _table_to_md(rows: list[list[str]]) -> str:
    if not rows:
        return "(empty)"
    max_cols = max(len(r) for r in rows)
    norm = [r + [""] * (max_cols - len(r)) for r in rows]
    for r in norm:
        for i, cell in enumerate(r):
            r[i] = cell.replace("|", "\\|").replace("\n", " ")
    header = "| " + " | ".join(norm[0]) + " |"
    sep = "| " + " | ".join("---" for _ in norm[0]) + " |"
    body_lines = ["| " + " | ".join(row) + " |" for row in norm[1:]]
    return "\n".join([header, sep] + body_lines)


def _extract_pdf(file_bytes: bytes) -> str:
    import pymupdf
    import pymupdf4llm

    doc = pymupdf.open(stream=file_bytes, filetype="pdf")
    md_text = pymupdf4llm.to_markdown(doc)
    doc.close()
    return f"{_TEXT_ONLY_NOTE}\n\n{md_text}"


def _get_docx_image_alt(shape) -> str:
    try:
        el = shape._element
        for child in el:
            if child.tag.endswith("}docPr") or child.tag == "docPr":
                return child.get("descr", "")
    except Exception:
        pass
    return ""


def _docx_runs_to_md(runs) -> str:
    parts = []
    for run in runs:
        text = run.text
        if not text:
            continue
        bold = run.bold
        italic = run.italic
        if bold and italic:
            text = f"***{text}***"
        elif bold:
            text = f"**{text}**"
        elif italic:
            text = f"*{text}*"
        parts.append(text)
    return "".join(parts)


def _docx_paragraph_to_md(p) -> str:
    style_name = (p.style.name or "").lower() if p.style else ""

    prefix = ""
    if style_name.startswith("heading"):
        try:
            level = int(style_name.split()[-1])
            prefix = "#" * min(level, 6) + " "
        except (ValueError, IndexError):
            prefix = "# "
    elif style_name == "title":
        prefix = "# "
    elif style_name == "subtitle":
        prefix = "## "
    elif style_name.startswith("list bullet"):
        prefix = "- "
    elif style_name.startswith("list number"):
        prefix = "1. "

    text = _docx_runs_to_md(p.runs)
    if not text.strip():
        return ""
    return prefix + text


def _extract_docx(file_bytes: bytes) -> str:
    from docx import Document
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(io.BytesIO(file_bytes))
    parts = []

    for child in doc.element.body:
        if child.tag == qn("w:p"):
            p = Paragraph(child, doc)
            md = _docx_paragraph_to_md(p)
            if md:
                parts.append(md)
        elif child.tag == qn("w:tbl"):
            tbl = Table(child, doc)
            rows = []
            for row in tbl.rows:
                rows.append([cell.text for cell in row.cells])
            md = _table_to_md(rows)
            parts.append(md)

    image_placeholders = []
    for shape in doc.inline_shapes:
        alt = _get_docx_image_alt(shape)
        image_placeholders.append(f"[IMAGE: {alt}]" if alt else "[IMAGE]")

    text = "\n\n".join(parts)
    if image_placeholders:
        text += "\n\n(Images found — original positions in document not preserved):\n" + "\n".join(
            image_placeholders
        )
    return f"{_TEXT_ONLY_NOTE}\n\n" + text


def _extract_xlsx(file_bytes: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            rows.append(cells)
        if rows:
            table_md = _table_to_md(rows)
        else:
            table_md = "(empty)"
        sheets.append(f"## Sheet: {sheet_name}\n{table_md}")
    wb.close()
    return f"{_CELLS_ONLY_NOTE}\n\n" + "\n\n".join(sheets)


def _get_pptx_shape_alt(shape) -> str:
    try:
        return shape._element.nvPicPr.cNvPr.get("descr", "")
    except (AttributeError, KeyError):
        return ""


def _pptx_text_frame_to_md(text_frame) -> str:
    lines = []
    for para in text_frame.paragraphs:
        parts = []
        for run in para.runs:
            text = run.text
            if not text:
                continue
            bold = run.font.bold
            italic = run.font.italic
            if bold and italic:
                text = f"***{text}***"
            elif bold:
                text = f"**{text}**"
            elif italic:
                text = f"*{text}*"
            parts.append(text)
        line = "".join(parts)
        if not line.strip():
            continue
        level = para.level or 0
        if level > 0:
            line = "  " * level + "- " + line
        lines.append(line)
    return "\n".join(lines)


def _extract_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        if slide.shapes.title and slide.shapes.title.text.strip():
            title = slide.shapes.title.text.strip()

        body_parts = []
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                alt = _get_pptx_shape_alt(shape)
                body_parts.append(f"[IMAGE: {alt}]" if alt else "[IMAGE]")
            elif shape.has_table:
                rows = []
                for row in shape.table.rows:
                    rows.append([cell.text for cell in row.cells])
                body_parts.append(_table_to_md(rows))
            elif shape.has_text_frame:
                text = _pptx_text_frame_to_md(shape.text_frame)
                if text:
                    body_parts.append(text)

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        header = f"## Slide {i}: {title}" if title else f"## Slide {i}"
        parts = [header]
        if body_parts:
            parts.append("\n".join(body_parts))
        if notes:
            parts.append(f"Notes: {notes}")
        slides.append("\n".join(parts))

    return f"{_TEXT_ONLY_NOTE}\n\n" + "\n\n".join(slides)


# ---------------------------------------------------------------------------
# Handler wrappers — return content block dicts
# ---------------------------------------------------------------------------


def _import_pdf(f: InputFile) -> ContentBlock:
    try:
        text = _extract_pdf(f["file_bytes"])
        return {"type": "text", "text": f"[File: {f['filename']}]\n\n{text}"}
    except Exception as exc:
        raise FileImportToolException(f"Failed to extract text from {f['filename']}") from exc


def _import_docx(f: InputFile) -> ContentBlock:
    try:
        text = _extract_docx(f["file_bytes"])
        return {"type": "text", "text": f"[File: {f['filename']}]\n\n{text}"}
    except Exception as exc:
        raise FileImportToolException(f"Failed to extract text from {f['filename']}") from exc


def _import_xlsx(f: InputFile) -> ContentBlock:
    try:
        text = _extract_xlsx(f["file_bytes"])
        return {"type": "text", "text": f"[File: {f['filename']}]\n\n{text}"}
    except Exception as exc:
        raise FileImportToolException(f"Failed to extract text from {f['filename']}") from exc


def _import_pptx(f: InputFile) -> ContentBlock:
    try:
        text = _extract_pptx(f["file_bytes"])
        return {"type": "text", "text": f"[File: {f['filename']}]\n\n{text}"}
    except Exception as exc:
        raise FileImportToolException(f"Failed to extract text from {f['filename']}") from exc


def _import_text(f: InputFile) -> ContentBlock:
    try:
        text = f["file_bytes"].decode("utf-8", errors="replace")
        return {"type": "text", "text": f"[File: {f['filename']}]\n\n{text}"}
    except Exception as exc:
        raise FileImportToolException(f"Failed to extract text from {f['filename']}") from exc


def _import_image(f: InputFile) -> ContentBlock:
    try:
        b64 = base64.standard_b64encode(f["file_bytes"]).decode("utf-8")
        return {
            "type": "image",
            "source": {
                "type": "base64",
                "media_type": f["mimetype"],
                "data": b64,
            },
        }
    except Exception as exc:
        raise FileImportToolException(f"Failed to process image {f['filename']}") from exc


_HANDLER_MANIFEST = [
    {
        "name": "import_pdf",
        "mimes": {"application/pdf"},
        "max_size": 10_000_000,
        "handler": _import_pdf,
    },
    {
        "name": "import_docx",
        "mimes": {MIME_DOCX},
        "max_size": 10_000_000,
        "handler": _import_docx,
    },
    {
        "name": "import_xlsx",
        "mimes": {MIME_XLSX},
        "max_size": 10_000_000,
        "handler": _import_xlsx,
    },
    {
        "name": "import_pptx",
        "mimes": {MIME_PPTX},
        "max_size": 10_000_000,
        "handler": _import_pptx,
    },
    {
        "name": "import_text",
        "mimes": {"text/plain", "text/csv", "text/markdown"},
        "max_size": 10_000_000,
        "handler": _import_text,
    },
    {
        "name": "import_image",
        "mimes": {"image/png", "image/jpeg", "image/gif", "image/webp"},
        "max_size": 10_000_000,
        "handler": _import_image,
    },
]


class Provider(BaseFileImporterProvider):
    """Built-in document import handlers (PDF, DOCX, XLSX, PPTX, text, images)."""

    def __init__(self, allowed_functions: list[str], **kwargs):
        super().__init__(allowed_functions, **kwargs)
        self._handler_map = {h["name"]: h["handler"] for h in _HANDLER_MANIFEST}

    def _get_all_tools(self) -> list[dict]:
        return _HANDLER_MANIFEST

    async def call_tool(
        self,
        name: str,
        arguments: dict,
        user_conversation_context: UserConversationContext,
        storage: BaseStorageProvider,
    ) -> ContentBlock:
        handler = self._handler_map.get(name)
        if not handler:
            raise FileImportToolException(f"Unknown import handler: {name}")
        return handler(arguments)
