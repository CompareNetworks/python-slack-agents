"""Built-in tool: document generation (PDF, DOCX, XLSX, CSV, PPTX).

Exports a Provider class that subclasses BaseToolProvider.
"""

import csv
import io
import logging
import re
from pathlib import Path

from docx import Document
from fpdf import FPDF
from openpyxl import Workbook

from slack_agents import UserConversationContext
from slack_agents.storage.base import BaseStorageProvider
from slack_agents.tools.base import BaseToolProvider, ToolResult

logger = logging.getLogger(__name__)

_FONT_DIR = Path(__file__).resolve().parent.parent.parent.parent / "fonts"
_FONT_REGULAR = _FONT_DIR / "DejaVuSans.ttf"
_FONT_BOLD = _FONT_DIR / "DejaVuSans-Bold.ttf"


def _sanitize_filename(name: str, ext: str) -> str:
    """Sanitize a string for use as a filename."""
    name = re.sub(r"[^\w\s-]", "", name).strip()
    name = re.sub(r"\s+", "_", name)
    return f"{name or 'document'}.{ext}"


def _parse_body_lines(body: str):
    """Parse markdown-ish body text into structured elements.

    Yields (type, data) tuples:
      ("h1", "Heading text")
      ("h2", "Sub-heading text")
      ("h3", "Sub-sub-heading text")
      ("bullet", "List item text")
      ("numbered", "Numbered item text")
      ("table", [["cell", ...], ...])
      ("paragraph", "Normal text")
    """
    _numbered_re = re.compile(r"^\d+\.\s")
    _separator_re = re.compile(r"^:?-+:?$")

    def _is_table_line(line: str) -> bool:
        s = line.strip()
        return s.startswith("|") and s.endswith("|")

    def _parse_table_line(line: str) -> list[str]:
        return [cell.strip() for cell in line.strip().strip("|").split("|")]

    def _is_separator(cells: list[str]) -> bool:
        return all(_separator_re.match(c) for c in cells)

    table_buf: list[list[str]] = []

    for line in body.split("\n"):
        stripped = line.strip()

        if _is_table_line(stripped):
            cells = _parse_table_line(stripped)
            if not _is_separator(cells):
                table_buf.append(cells)
            continue

        if table_buf:
            yield ("table", table_buf)
            table_buf = []

        if not stripped:
            yield ("paragraph", "")
        elif stripped.startswith("### "):
            yield ("h3", stripped[4:])
        elif stripped.startswith("## "):
            yield ("h2", stripped[3:])
        elif stripped.startswith("# "):
            yield ("h1", stripped[2:])
        elif stripped.startswith("- ") or stripped.startswith("* "):
            yield ("bullet", stripped[2:])
        elif _numbered_re.match(stripped):
            yield ("numbered", _numbered_re.sub("", stripped, count=1))
        else:
            yield ("paragraph", stripped)

    if table_buf:
        yield ("table", table_buf)


def _make_pdf() -> FPDF:
    """Create an FPDF instance with Unicode font support."""
    pdf = FPDF()
    if _FONT_REGULAR.exists() and _FONT_BOLD.exists():
        pdf.add_font("DejaVu", "", str(_FONT_REGULAR))
        pdf.add_font("DejaVu", "B", str(_FONT_BOLD))
        pdf._font_family = "DejaVu"  # noqa: SLF001
    else:
        logger.warning(
            "DejaVu fonts not found at %s — falling back to Helvetica (latin-1 only). "
            "Run: python -m slack_agents.scripts.download_fonts",
            _FONT_DIR,
        )
        pdf._font_family = "Helvetica"  # noqa: SLF001
    return pdf


def _pdf_font(pdf: FPDF) -> str:
    return pdf._font_family  # noqa: SLF001


_RICH_TEXT_RE = re.compile(r"(\*\*.*?\*\*|\*(?!\*).*?(?<!\*)\*(?!\*))")


def _write_rich_text(pdf: FPDF, text: str, line_height: float) -> None:
    """Write text with **bold** and *italic* support to a PDF."""
    font = _pdf_font(pdf)
    is_ttf = font != "Helvetica"
    parts = _RICH_TEXT_RE.split(text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            pdf.set_font(font, "B", 11)
            pdf.write(line_height, part[2:-2])
            pdf.set_font(font, "", 11)
        elif part.startswith("*") and part.endswith("*") and len(part) > 2:
            if not is_ttf:
                pdf.set_font(font, "I", 11)
            pdf.write(line_height, part[1:-1])
            if not is_ttf:
                pdf.set_font(font, "", 11)
        else:
            pdf.write(line_height, part)


async def _export_pdf(arguments: dict) -> dict:
    title = arguments["title"]
    body = arguments["body"]
    filename = _sanitize_filename(title, "pdf")

    pdf = _make_pdf()
    font = _pdf_font(pdf)
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    pdf.set_font(font, "B", 18)
    pdf.cell(0, 12, title, new_x="LMARGIN", new_y="NEXT")
    pdf.ln(4)

    pdf.set_font(font, "", 11)
    is_ttf = font != "Helvetica"
    bullet_char = "\u2022" if is_ttf else "-"
    line_height = 6

    for elem_type, data in _parse_body_lines(body):
        if elem_type == "table":
            pdf.ln(2)
            with pdf.table(first_row_as_headings=True) as table:
                for row_cells in data:
                    row = table.row()
                    for cell in row_cells:
                        row.cell(cell)
            pdf.ln(2)
            pdf.set_font(font, "", 11)
        elif elem_type == "h1":
            pdf.ln(4)
            pdf.set_font(font, "B", 15)
            pdf.cell(0, 8, data, new_x="LMARGIN", new_y="NEXT")
            pdf.set_font(font, "", 11)
        elif elem_type == "h2":
            pdf.ln(2)
            pdf.set_font(font, "B", 13)
            pdf.cell(0, 7, data, new_x="LMARGIN", new_y="NEXT")
            pdf.set_font(font, "", 11)
        elif elem_type == "h3":
            pdf.ln(2)
            pdf.set_font(font, "B", 12)
            pdf.cell(0, 7, data, new_x="LMARGIN", new_y="NEXT")
            pdf.set_font(font, "", 11)
        elif elem_type == "bullet":
            pdf.cell(6)
            _write_rich_text(pdf, f"{bullet_char} {data}", line_height)
            pdf.ln(line_height)
        elif elem_type == "numbered":
            pdf.cell(6)
            _write_rich_text(pdf, data, line_height)
            pdf.ln(line_height)
        elif data:
            _write_rich_text(pdf, data, line_height)
            pdf.ln(line_height)
        else:
            pdf.ln(line_height // 2)

    data = pdf.output()
    pages = pdf.pages_count
    return {
        "content": f"Created {filename} ({pages} page{'s' if pages != 1 else ''})",
        "is_error": False,
        "files": [{"data": bytes(data), "filename": filename, "mimeType": "application/pdf"}],
    }


def _add_rich_paragraph(doc: Document, text: str) -> None:
    para = doc.add_paragraph()
    _add_rich_runs(para, text)


def _add_rich_runs(para, text: str) -> None:
    parts = _RICH_TEXT_RE.split(text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            para.add_run(part[2:-2]).bold = True
        elif part.startswith("*") and part.endswith("*") and len(part) > 2:
            para.add_run(part[1:-1]).italic = True
        else:
            para.add_run(part)


async def _export_docx(arguments: dict) -> dict:
    title = arguments["title"]
    body = arguments["body"]
    filename = _sanitize_filename(title, "docx")

    doc = Document()
    doc.add_heading(title, level=0)

    for elem_type, data in _parse_body_lines(body):
        if elem_type == "table":
            if data:
                tbl = doc.add_table(rows=len(data), cols=len(data[0]), style="Table Grid")
                for i, row_cells in enumerate(data):
                    for j, cell_text in enumerate(row_cells):
                        tbl.rows[i].cells[j].text = cell_text
        elif elem_type == "h1":
            doc.add_heading(data, level=1)
        elif elem_type == "h2":
            doc.add_heading(data, level=2)
        elif elem_type == "h3":
            doc.add_heading(data, level=3)
        elif elem_type == "bullet":
            doc.add_paragraph(data, style="List Bullet")
        elif elem_type == "numbered":
            doc.add_paragraph(data, style="List Number")
        elif data:
            _add_rich_paragraph(doc, data)

    buf = io.BytesIO()
    doc.save(buf)
    buf_data = buf.getvalue()

    return {
        "content": f"Created {filename}",
        "is_error": False,
        "files": [
            {
                "data": buf_data,
                "filename": filename,
                "mimeType": (
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                ),
            }
        ],
    }


async def _export_xlsx(arguments: dict) -> dict:
    filename_base = arguments["filename"]
    sheets = arguments["sheets"]
    filename = _sanitize_filename(filename_base, "xlsx")

    wb = Workbook()
    wb.remove(wb.active)

    total_rows = 0
    for sheet_def in sheets:
        ws = wb.create_sheet(title=sheet_def["name"])
        for row in sheet_def["rows"]:
            ws.append(row)
            total_rows += 1

    buf = io.BytesIO()
    wb.save(buf)
    data = buf.getvalue()

    sheet_count = len(sheets)
    return {
        "content": (
            f"Created {filename} ({sheet_count} sheet{'s' if sheet_count != 1 else ''}, "
            f"{total_rows} rows)"
        ),
        "is_error": False,
        "files": [
            {
                "data": data,
                "filename": filename,
                "mimeType": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            }
        ],
    }


async def _export_csv(arguments: dict) -> dict:
    filename_base = arguments["filename"]
    rows = arguments["rows"]
    filename = _sanitize_filename(filename_base, "csv")

    buf = io.StringIO()
    writer = csv.writer(buf)
    for row in rows:
        writer.writerow(row)
    data = buf.getvalue().encode("utf-8")

    return {
        "content": f"Created {filename} ({len(rows)} rows)",
        "is_error": False,
        "files": [{"data": data, "filename": filename, "mimeType": "text/csv"}],
    }


def _pptx_add_rich_text(paragraph, text: str) -> None:
    parts = _RICH_TEXT_RE.split(text)
    for part in parts:
        run = paragraph.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]
            run.font.bold = True
        elif part.startswith("*") and part.endswith("*") and len(part) > 2:
            run.text = part[1:-1]
            run.font.italic = True
        else:
            run.text = part


async def _export_pptx(arguments: dict) -> dict:
    from pptx import Presentation
    from pptx.util import Emu, Inches

    title = arguments["title"]
    slides_data = arguments["slides"]
    filename = _sanitize_filename(title, "pptx")

    prs = Presentation()
    for slide_def in slides_data:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = slide_def["title"]

        body_placeholder = slide.placeholders[1]
        tf = body_placeholder.text_frame
        tf.clear()
        first = True

        elements = list(_parse_body_lines(slide_def.get("body", "")))
        for elem_type, data in elements:
            if elem_type == "table":
                if data:
                    rows_count = len(data)
                    cols_count = len(data[0])
                    ph_left = body_placeholder.left
                    ph_top = body_placeholder.top + body_placeholder.height + Inches(0.1)
                    ph_width = body_placeholder.width
                    row_height = Emu(370000)
                    tbl_height = row_height * rows_count
                    shape = slide.shapes.add_table(
                        rows_count, cols_count, ph_left, ph_top, ph_width, tbl_height
                    )
                    for i, row_cells in enumerate(data):
                        for j, cell_text in enumerate(row_cells):
                            shape.table.cell(i, j).text = cell_text
                continue

            if elem_type == "paragraph" and not data:
                continue

            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()

            if elem_type == "bullet":
                p.level = 1

            _pptx_add_rich_text(p, data)

    buf = io.BytesIO()
    prs.save(buf)
    buf_data = buf.getvalue()

    slide_count = len(slides_data)
    return {
        "content": (f"Created {filename} ({slide_count} slide{'s' if slide_count != 1 else ''})"),
        "is_error": False,
        "files": [
            {
                "data": buf_data,
                "filename": filename,
                "mimeType": (
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                ),
            }
        ],
    }


_TOOL_MANIFEST = [
    {
        "name": "export_pdf",
        "description": (
            "Generate a simple PDF document with a basic layout meant as a starting point. "
            "The body supports formatting: #, ##, ### for headings, **bold**, *italic*, "
            "lines starting with - for bullet lists, numbered lists (1. ), and markdown "
            "tables (| col1 | col2 | with |---| separator). "
            "No images, charts, or advanced styling."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "Document title (used as filename)"},
                "body": {
                    "type": "string",
                    "description": "Document body with optional markdown-style formatting",
                },
            },
            "required": ["title", "body"],
        },
        "handler": _export_pdf,
    },
    {
        "name": "export_docx",
        "description": (
            "Generate a simple Word (.docx) document with a basic layout meant as a starting "
            "point. The body supports formatting: #, ##, ### for headings, **bold**, *italic*, "
            "lines starting with - for bullet lists, numbered lists (1. ), and markdown "
            "tables (| col1 | col2 | with |---| separator). "
            "No images, charts, or advanced styling."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "Document title (used as filename)"},
                "body": {
                    "type": "string",
                    "description": "Document body with optional markdown-style formatting",
                },
            },
            "required": ["title", "body"],
        },
        "handler": _export_docx,
    },
    {
        "name": "export_xlsx",
        "description": (
            "Generate a simple Excel (.xlsx) spreadsheet from structured data. "
            "Produces a basic layout meant as a starting point — text and numbers only, "
            "no formulas, charts, or styling."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "filename": {
                    "type": "string",
                    "description": "Filename for the spreadsheet (without extension)",
                },
                "sheets": {
                    "type": "array",
                    "description": "List of sheets, each with a name and rows of data",
                    "items": {
                        "type": "object",
                        "properties": {
                            "name": {"type": "string", "description": "Sheet name"},
                            "rows": {
                                "type": "array",
                                "description": "Rows of data (first row is typically headers)",
                                "items": {
                                    "type": "array",
                                    "items": {"type": "string"},
                                },
                            },
                        },
                        "required": ["name", "rows"],
                    },
                },
            },
            "required": ["filename", "sheets"],
        },
        "handler": _export_xlsx,
    },
    {
        "name": "export_csv",
        "description": "Generate a simple CSV file from rows of data. Plain text only.",
        "input_schema": {
            "type": "object",
            "properties": {
                "filename": {
                    "type": "string",
                    "description": "Filename for the CSV (without extension)",
                },
                "rows": {
                    "type": "array",
                    "description": "Rows of data (first row is typically headers)",
                    "items": {
                        "type": "array",
                        "items": {"type": "string"},
                    },
                },
            },
            "required": ["filename", "rows"],
        },
        "handler": _export_csv,
    },
    {
        "name": "export_pptx",
        "description": (
            "Generate a simple PowerPoint (.pptx) presentation with a basic layout meant as "
            "a starting point. Provide a title and an array of slides, each with a title and "
            "body. Body supports **bold**, *italic*, lines starting with - for bullet lists, "
            "numbered lists (1. ), #/##/### headings, and markdown tables "
            "(| col1 | col2 | with |---| separator). "
            "No images, charts, or advanced styling."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {
                    "type": "string",
                    "description": "Presentation title (used as filename)",
                },
                "slides": {
                    "type": "array",
                    "description": "List of slides, each with a title and body",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string", "description": "Slide title"},
                            "body": {
                                "type": "string",
                                "description": (
                                    "Slide body text; lines starting with - become bullet points"
                                ),
                            },
                        },
                        "required": ["title"],
                    },
                },
            },
            "required": ["title", "slides"],
        },
        "handler": _export_pptx,
    },
]


class Provider(BaseToolProvider):
    """Built-in document export tools (PDF, DOCX, XLSX, CSV, PPTX)."""

    def __init__(self, allowed_functions: list[str]):
        super().__init__(allowed_functions)
        self._handlers = {t["name"]: t["handler"] for t in _TOOL_MANIFEST}

    def _get_all_tools(self) -> list[dict]:
        return [
            {"name": t["name"], "description": t["description"], "input_schema": t["input_schema"]}
            for t in _TOOL_MANIFEST
        ]

    async def call_tool(
        self,
        name: str,
        arguments: dict,
        user_conversation_context: UserConversationContext,
        storage: BaseStorageProvider,
    ) -> ToolResult:
        handler = self._handlers.get(name)
        if not handler:
            return {"content": f"Unknown tool: {name}", "is_error": True, "files": []}
        try:
            return await handler(arguments)
        except Exception as e:
            logger.exception("Export tool call failed: %s", name)
            return {"content": f"Tool execution error: {e}", "is_error": True, "files": []}
